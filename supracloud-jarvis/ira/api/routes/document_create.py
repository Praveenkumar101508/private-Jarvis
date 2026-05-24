"""
IRA Native Document Generation — Feature #4.

Create professional PDF, Word, PowerPoint, and Excel files directly from chat.

POST /document/create  — generate a document from a prompt
GET  /document/download/{doc_id} — download a generated document

Trigger phrases:
  "create a PDF of...", "make a Word document with...", "generate a PowerPoint...",
  "write an Excel spreadsheet...", "build a report on...", "create a slide deck..."

Libraries:
  PDF:  reportlab (pip install reportlab)
  DOCX: python-docx (already in requirements)
  PPTX: python-pptx (pip install python-pptx)
  XLSX: openpyxl (pip install openpyxl)
"""

from __future__ import annotations

import base64
import io
import json as _json
import logging
import re
import time
import uuid
from pathlib import Path
from typing import Literal, Optional

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field
from sse_starlette.sse import EventSourceResponse

from api.middleware.auth import require_auth
from utils.llm import chat_complete

router = APIRouter(prefix="/document", tags=["document"])
logger = logging.getLogger("ira.document_create")

_DOC_TTL = 3600  # 1 hour — documents expire after download window

# ── Redis-backed document store (per-user, TTL-protected) ────────────────────
# Falls back to a process-local dict when Redis is unavailable (e.g., dev mode).
_DOC_FALLBACK: dict[str, dict] = {}


async def _store_doc(user_id: str, doc_id: str, file_bytes: bytes, filename: str, mimetype: str) -> None:
    payload = _json.dumps({
        "data": base64.b64encode(file_bytes).decode(),
        "filename": filename,
        "mimetype": mimetype,
        "owner": user_id,
    })
    try:
        from utils.redis_client import get_redis
        redis = get_redis()
        await redis.setex(f"document:{user_id}:{doc_id}", _DOC_TTL, payload)
    except Exception:
        _DOC_FALLBACK[f"{user_id}:{doc_id}"] = {
            "data": file_bytes, "filename": filename, "mimetype": mimetype, "owner": user_id,
        }


async def _get_doc(user_id: str, doc_id: str) -> tuple[bytes, str, str] | None:
    """Return (bytes, filename, mimetype) for this user's doc, or None if not found/expired."""
    try:
        from utils.redis_client import get_redis
        redis = get_redis()
        raw = await redis.get(f"document:{user_id}:{doc_id}")
        if not raw:
            return None
        entry = _json.loads(raw)
        if entry.get("owner") != user_id:
            return None  # Cross-user access blocked
        return base64.b64decode(entry["data"]), entry["filename"], entry["mimetype"]
    except Exception:
        fb = _DOC_FALLBACK.get(f"{user_id}:{doc_id}")
        if fb and fb.get("owner") == user_id:
            return fb["data"], fb["filename"], fb["mimetype"]
        return None

# Document type detection
_DOC_CREATE_RE = re.compile(
    r"\b(create\s+(a\s+)?(pdf|word|docx|powerpoint|pptx|excel|xlsx|spreadsheet|"
    r"slide.?deck|presentation|report|document|memo|proposal|invoice|contract)|"
    r"make\s+(a\s+)?(pdf|word|powerpoint|excel|report|document|slide)|"
    r"generate\s+(a\s+)?(pdf|word|report|presentation|document|invoice)|"
    r"write\s+(a\s+)?(report|document|memo|proposal|contract)|"
    r"build\s+(a\s+)?(report|presentation|spreadsheet|document))\b",
    re.I,
)

_FORMAT_RE = re.compile(
    r"\b(pdf|word|docx|powerpoint|pptx|excel|xlsx|spreadsheet|presentation|slides?)\b",
    re.I,
)


def is_doc_create_request(query: str) -> bool:
    return bool(_DOC_CREATE_RE.search(query))


def _detect_format(query: str) -> Literal["pdf", "docx", "pptx", "xlsx"]:
    m = _FORMAT_RE.search(query)
    if not m:
        return "pdf"
    t = m.group(1).lower()
    if t in ("excel", "xlsx", "spreadsheet"):
        return "xlsx"
    if t in ("powerpoint", "pptx", "presentation", "slides", "slide"):
        return "pptx"
    if t in ("word", "docx"):
        return "docx"
    return "pdf"


# ── Request model ─────────────────────────────────────────────────────────────

class DocCreateRequest(BaseModel):
    prompt: str = Field(..., min_length=1, max_length=8000)
    format: Optional[Literal["pdf", "docx", "pptx", "xlsx"]] = None
    session_id: str = Field(default_factory=lambda: str(uuid.uuid4()))


# ── LLM content generation ────────────────────────────────────────────────────

_DOC_SYSTEM = """\
You are a professional document writer. Generate structured document content based on the user's request.

Output the document content in clean Markdown with:
- Clear headings (# ## ###)
- Bullet points and numbered lists where appropriate
- Tables using Markdown table syntax where relevant
- Bold key terms
- Professional, concise language

For presentations: structure as ## Slide 1: Title / content / ## Slide 2: ...
For spreadsheets: structure as a Markdown table with headers and data rows
For reports/docs: standard professional document structure

Generate COMPLETE content — do not truncate or say "continue...".
"""


async def _generate_doc_content(prompt: str, fmt: str) -> str:
    fmt_hints = {
        "pdf": "Create a professional PDF report/document.",
        "docx": "Create a professional Word document.",
        "pptx": "Create a presentation with clear slide titles and content. Format each slide as '## Slide N: [Title]' followed by content.",
        "xlsx": "Create a spreadsheet. Format data as a proper Markdown table with headers. Include multiple sections if needed.",
    }
    msgs = [
        {"role": "system", "content": _DOC_SYSTEM},
        {"role": "user", "content": f"{fmt_hints.get(fmt, '')}\n\n{prompt}"},
    ]
    return await chat_complete(msgs, use_deep=True, max_tokens=4096, temperature=0.4)


# ── PDF generation ────────────────────────────────────────────────────────────

def _generate_pdf(content: str, title: str) -> bytes:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib import colors

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4, rightMargin=25*mm, leftMargin=25*mm,
                                topMargin=25*mm, bottomMargin=25*mm)
        styles = getSampleStyleSheet()
        story = []

        # Title
        title_style = ParagraphStyle("CustomTitle", parent=styles["Title"], fontSize=18, spaceAfter=16)
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 8))

        # Parse markdown-ish content into reportlab elements
        for line in content.split("\n"):
            line = line.strip()
            if not line:
                story.append(Spacer(1, 6))
                continue
            if line.startswith("### "):
                story.append(Paragraph(line[4:], styles["Heading3"]))
            elif line.startswith("## "):
                story.append(Paragraph(line[3:], styles["Heading2"]))
            elif line.startswith("# "):
                story.append(Paragraph(line[2:], styles["Heading1"]))
            elif line.startswith(("- ", "* ", "• ")):
                text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', line[2:])
                story.append(Paragraph(f"• {text}", styles["Normal"]))
            elif re.match(r"^\d+\. ", line):
                text = re.sub(r"^\d+\. ", "", line)
                story.append(Paragraph(text, styles["Normal"]))
            else:
                # Bold markers
                line = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", line)
                story.append(Paragraph(line, styles["Normal"]))

        doc.build(story)
        return buf.getvalue()
    except ImportError:
        # Fallback: plain text PDF via basic bytes
        logger.warning("reportlab not installed — generating plain text fallback")
        text = f"{title}\n{'='*len(title)}\n\n{content}"
        return text.encode("utf-8")


# ── DOCX generation ───────────────────────────────────────────────────────────

def _generate_docx(content: str, title: str) -> bytes:
    try:
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()
        doc.add_heading(title, 0)

        for line in content.split("\n"):
            line = line.strip()
            if not line:
                doc.add_paragraph()
                continue
            if line.startswith("### "):
                doc.add_heading(line[4:], level=3)
            elif line.startswith("## "):
                doc.add_heading(line[3:], level=2)
            elif line.startswith("# "):
                doc.add_heading(line[2:], level=1)
            elif line.startswith(("- ", "* ", "• ")):
                p = doc.add_paragraph(line[2:], style="List Bullet")
            elif re.match(r"^\d+\. ", line):
                p = doc.add_paragraph(re.sub(r"^\d+\. ", "", line), style="List Number")
            else:
                p = doc.add_paragraph()
                # Handle **bold** inline
                parts = re.split(r"\*\*(.+?)\*\*", line)
                for i, part in enumerate(parts):
                    run = p.add_run(part)
                    if i % 2 == 1:
                        run.bold = True

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()
    except ImportError:
        logger.warning("python-docx not installed")
        return content.encode("utf-8")


# ── PPTX generation ───────────────────────────────────────────────────────────

def _generate_pptx(content: str, title: str) -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if slide.placeholders[1].has_text_frame:
            slide.placeholders[1].text = "Generated by IRA"

        # Parse slide markers
        current_slide = None
        current_bullets: list[str] = []

        def _flush_slide(stitle: str, bullets: list[str]) -> None:
            layout = prs.slide_layouts[1]  # Title + Content
            sl = prs.slides.add_slide(layout)
            sl.shapes.title.text = stitle
            tf = sl.placeholders[1].text_frame
            tf.clear()
            for i, b in enumerate(bullets[:8]):
                if i == 0:
                    tf.paragraphs[0].text = b
                else:
                    p = tf.add_paragraph()
                    p.text = b
                    p.level = 1

        for line in content.split("\n"):
            line = line.strip()
            if line.startswith("## Slide") or line.startswith("##"):
                if current_slide and current_bullets:
                    _flush_slide(current_slide, current_bullets)
                current_slide = re.sub(r"^##\s*(Slide\s*\d*:?\s*)?", "", line).strip()
                current_bullets = []
            elif line.startswith(("- ", "* ", "• ")):
                current_bullets.append(line[2:])
            elif line and not line.startswith("#") and current_slide:
                current_bullets.append(line)

        if current_slide:
            _flush_slide(current_slide, current_bullets)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    except ImportError:
        logger.warning("python-pptx not installed — pip install python-pptx")
        return content.encode("utf-8")


# ── XLSX generation ───────────────────────────────────────────────────────────

def _generate_xlsx(content: str, title: str) -> bytes:
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = title[:31]  # Excel sheet name limit

        # Parse markdown tables from content
        row_num = 1
        ws.cell(row=row_num, column=1, value=title).font = Font(bold=True, size=14)
        row_num += 2

        in_table = False
        for line in content.split("\n"):
            line = line.strip()
            if not line:
                if in_table:
                    row_num += 1
                    in_table = False
                continue
            if "|" in line:
                cells = [c.strip() for c in line.split("|") if c.strip()]
                if not cells:
                    continue
                if re.match(r"^[-:\s|]+$", line):
                    # Table separator row — make previous row bold header
                    for col_idx in range(1, ws.max_column + 1):
                        cell = ws.cell(row=row_num - 1, column=col_idx)
                        cell.font = Font(bold=True)
                        cell.fill = PatternFill("solid", fgColor="D3D3D3")
                    continue
                for col_idx, val in enumerate(cells, 1):
                    ws.cell(row=row_num, column=col_idx, value=val)
                row_num += 1
                in_table = True
            elif line.startswith("#"):
                row_num += 1
                cell = ws.cell(row=row_num, column=1, value=line.lstrip("#").strip())
                cell.font = Font(bold=True, size=12)
                row_num += 1
            else:
                ws.cell(row=row_num, column=1, value=line)
                row_num += 1

        # Auto-size columns
        for col in ws.columns:
            max_len = max((len(str(cell.value or "")) for cell in col), default=10)
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 60)

        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()
    except ImportError:
        logger.warning("openpyxl not installed — pip install openpyxl")
        return content.encode("utf-8")


# ── MIME types ────────────────────────────────────────────────────────────────
_MIME = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}
_EXT = {"pdf": ".pdf", "docx": ".docx", "pptx": ".pptx", "xlsx": ".xlsx"}


# ── SSE create endpoint ───────────────────────────────────────────────────────

@router.post("/create")
async def document_create(
    req: DocCreateRequest,
    _user: str = Depends(require_auth),
):
    """Generate a document from a prompt and return a download link (SSE)."""
    fmt = req.format or _detect_format(req.prompt)

    async def gen():
        t0 = time.monotonic()
        doc_name = re.sub(r"[^a-z0-9_]", "_", req.prompt[:40].lower()).strip("_")
        filename = f"ira_{doc_name}{_EXT[fmt]}"

        yield {"data": _json.dumps({"token": f"📄 Generating {fmt.upper()} document…\n\n"})}

        try:
            # Step 1: Generate content with LLM
            content = await _generate_doc_content(req.prompt, fmt)
            yield {"data": _json.dumps({"token": "✅ Content generated — building document…\n"})}

            # Step 2: Convert to file bytes
            title = doc_name.replace("_", " ").title()
            if fmt == "pdf":
                file_bytes = _generate_pdf(content, title)
            elif fmt == "docx":
                file_bytes = _generate_docx(content, title)
            elif fmt == "pptx":
                file_bytes = _generate_pptx(content, title)
            else:  # xlsx
                file_bytes = _generate_xlsx(content, title)

            # Step 3: Store with full UUID (not truncated 8-char) and per-user key
            doc_id = str(uuid.uuid4())
            await _store_doc(_user, doc_id, file_bytes, filename, _MIME[fmt])

            latency = int((time.monotonic() - t0) * 1000)
            yield {"data": _json.dumps({
                "document_created": True,
                "doc_id": doc_id,
                "filename": filename,
                "format": fmt,
                "size_kb": len(file_bytes) // 1024,
                "download_url": f"/api/v1/document/download/{doc_id}",
            })}
            yield {"data": _json.dumps({
                "token": (
                    f"\n📎 **{filename}** ready to download!\n"
                    f"Size: {len(file_bytes)//1024}KB · [{fmt.upper()}]\n\n"
                    f"[⬇️ Download]({f'/api/v1/document/download/{doc_id}'})"
                )
            })}
        except Exception as e:
            yield {"data": _json.dumps({"token": f"\n❌ Document generation error: {str(e)[:200]}"})}

        yield {"data": _json.dumps({
            "done": True, "agent": "document_create",
            "latency_ms": int((time.monotonic() - t0) * 1000),
        })}

    return EventSourceResponse(gen())


# ── Download endpoint ─────────────────────────────────────────────────────────

@router.get("/download/{doc_id}")
async def document_download(
    doc_id: str,
    _user: str = Depends(require_auth),
):
    """Download a previously generated document by its ID (owner-only, 1 h TTL)."""
    entry = await _get_doc(_user, doc_id)
    if not entry:
        raise HTTPException(status_code=404, detail="Document not found or expired.")
    file_bytes, filename, mimetype = entry
    return StreamingResponse(
        io.BytesIO(file_bytes),
        media_type=mimetype,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
